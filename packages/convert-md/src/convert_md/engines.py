"""Conversion engines (ADR 0021 / R142).

Each engine wraps one library, normalizes the output to Markdown, and grades
fidelity via :mod:`convert_md.fidelity`. Heavy imports (docling/torch,
pandoc, markitdown) live inside ``convert`` so importing this module is cheap —
only the engine actually selected pays the import cost.

Engines raise :class:`ConversionError` on failure so the dispatcher can walk
the format's fallback chain.
"""

from __future__ import annotations

from importlib.metadata import PackageNotFoundError
from importlib.metadata import version as _pkg_version
from typing import TYPE_CHECKING

from convert_md.base import ConversionError, ConversionResult
from convert_md.fidelity import grade

if TYPE_CHECKING:
    from pathlib import Path


def _ver(pkg: str) -> str:
    try:
        return _pkg_version(pkg)
    except PackageNotFoundError:
        return "unknown"


def _result(markdown: str, engine: str, source_size: int, *, check_yield: bool = True) -> ConversionResult:
    fidelity, lost, warnings = grade(markdown=markdown, source_size=source_size, check_yield=check_yield)
    return ConversionResult(body_markdown=markdown, engine=engine, fidelity=fidelity, lost=lost, warnings=warnings)


class PymupdfLlmEngine:
    """PDF → Markdown via pymupdf4llm (AGPL-3.0-via-PyMuPDF, no ML/GPU deps).

    Sole `.pdf` engine as of convert-md v0.5.0 (bench/results/2026-07-09-findings.md,
    design.md D5): docling was dropped entirely, not just demoted to fallback — the
    deciding evidence was a semantic-recoverability re-scoring (can an LLM recover
    the facts regardless of formatting), under which docling has a total,
    unrecoverable failure (drops 100% of formulas on math-heavy content) that this
    engine doesn't share. Also conflict-free: none of docling's unconditional
    ``typer<0.25`` pin, for every consumer, not just ones that can absorb the tax.
    """

    name = f"pymupdf4llm@{_ver('pymupdf4llm')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert a PDF to Markdown via pymupdf4llm."""
        try:
            import pymupdf4llm  # noqa: PLC0415 — lazy heavy import
        except ImportError as exc:  # pragma: no cover
            msg = f"pymupdf4llm unavailable: {exc}"
            raise ConversionError(msg) from exc
        try:
            markdown = pymupdf4llm.to_markdown(str(path))
        except Exception as exc:
            msg = f"pymupdf4llm failed on {path.name}: {exc}"
            raise ConversionError(msg) from exc
        return _result(markdown, self.name, path.stat().st_size, check_yield=False)


class PandocEngine:
    """DOCX → Markdown via pandoc, preserving tracked changes."""

    name = f"pandoc@{_ver('pypandoc-binary')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert a DOCX to GitHub-flavored Markdown via pandoc, keeping tracked changes."""
        try:
            import pypandoc  # noqa: PLC0415 — lazy heavy import
        except ImportError as exc:  # pragma: no cover
            msg = f"pypandoc unavailable: {exc}"
            raise ConversionError(msg) from exc
        try:
            engine = f"pandoc@{pypandoc.get_pandoc_version()}"
            markdown = pypandoc.convert_file(
                str(path),
                "gfm",
                extra_args=["--track-changes=all", "--wrap=none"],
            )
        except Exception as exc:
            msg = f"pandoc failed on {path.name}: {exc}"
            raise ConversionError(msg) from exc
        return _result(markdown, engine, path.stat().st_size, check_yield=False)


class MarkitdownEngine:
    """PPTX/XLSX/generic → Markdown via markitdown."""

    name = f"markitdown@{_ver('markitdown')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert a PPTX/XLSX/generic file to Markdown via markitdown, appending PPTX speaker notes."""
        try:
            from markitdown import MarkItDown  # noqa: PLC0415 — lazy heavy import
        except ImportError as exc:  # pragma: no cover
            msg = f"markitdown unavailable: {exc}"
            raise ConversionError(msg) from exc
        try:
            markdown = MarkItDown().convert(str(path)).text_content
        except Exception as exc:
            msg = f"markitdown failed on {path.name}: {exc}"
            raise ConversionError(msg) from exc
        if path.suffix.lower() == ".pptx":
            notes = _pptx_notes(path)
            if notes:
                markdown = f"{markdown}\n\n## Speaker notes\n\n{notes}"
        return _result(markdown, self.name, path.stat().st_size, check_yield=False)


class OpenpyxlEngine:
    """XLSX → Markdown tables via openpyxl (per-sheet, fallback for markitdown)."""

    name = f"openpyxl@{_ver('openpyxl')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert an XLSX to Markdown tables via openpyxl, one section per worksheet."""
        try:
            from openpyxl import load_workbook  # noqa: PLC0415 — lazy import
        except ImportError as exc:  # pragma: no cover
            msg = f"openpyxl unavailable: {exc}"
            raise ConversionError(msg) from exc
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                parts.append(f"## {ws.title}\n")
                rows = list(ws.iter_rows(values_only=True))
                parts.extend(_rows_to_markdown(rows))
            markdown = "\n".join(parts)
        except Exception as exc:
            msg = f"openpyxl failed on {path.name}: {exc}"
            raise ConversionError(msg) from exc
        return _result(markdown, self.name, path.stat().st_size, check_yield=False)


class TrafilaturaEngine:
    """HTML → Markdown via trafilatura (best boilerplate removal)."""

    name = f"trafilatura@{_ver('trafilatura')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert an HTML file to Markdown via trafilatura, stripping boilerplate."""
        try:
            import trafilatura  # noqa: PLC0415 — lazy import
        except ImportError as exc:  # pragma: no cover
            msg = f"trafilatura unavailable: {exc}"
            raise ConversionError(msg) from exc
        html = path.read_text(encoding="utf-8", errors="replace")
        markdown = trafilatura.extract(html, output_format="markdown", include_tables=True)
        if not markdown:
            msg = f"trafilatura extracted nothing from {path.name}"
            raise ConversionError(msg)
        return _result(markdown, self.name, path.stat().st_size)


class Html2TextEngine:
    """HTML → Markdown via html2text (fallback)."""

    name = f"html2text@{_ver('html2text')}"

    def convert(self, path: Path) -> ConversionResult:
        """Convert an HTML file to Markdown via html2text (unbounded line width)."""
        try:
            import html2text  # noqa: PLC0415 — lazy import
        except ImportError as exc:  # pragma: no cover
            msg = f"html2text unavailable: {exc}"
            raise ConversionError(msg) from exc
        html = path.read_text(encoding="utf-8", errors="replace")
        converter = html2text.HTML2Text()
        converter.body_width = 0
        markdown = converter.handle(html)
        return _result(markdown, self.name, path.stat().st_size)


# --- helpers -----------------------------------------------------------------


def _pptx_notes(path: Path) -> str:
    """Recover speaker notes markitdown drops (R142 pptx note pass)."""
    try:
        from pptx import Presentation  # noqa: PLC0415 — lazy import
    except ImportError:  # pragma: no cover
        return ""
    try:
        prs = Presentation(str(path))
    except Exception:  # noqa: BLE001  -- fallback chain: any engine failure falls through to the next engine
        return ""
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                chunks.append(f"**Slide {i}:** {text}")
    return "\n\n".join(chunks)


def _rows_to_markdown(rows: list[tuple]) -> list[str]:
    if not rows:
        return [""]
    out: list[str] = []
    header = rows[0]
    out.append("| " + " | ".join(_cell(c) for c in header) + " |")
    out.append("| " + " | ".join("---" for _ in header) + " |")
    out.extend("| " + " | ".join(_cell(c) for c in row) + " |" for row in rows[1:])
    out.append("")
    return out


def _cell(value: object) -> str:
    return "" if value is None else str(value).replace("|", "\\|").replace("\n", " ")


__all__ = [
    "Html2TextEngine",
    "MarkitdownEngine",
    "OpenpyxlEngine",
    "PandocEngine",
    "PymupdfLlmEngine",
    "TrafilaturaEngine",
]
